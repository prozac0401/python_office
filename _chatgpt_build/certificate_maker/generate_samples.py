# -*- coding: utf-8 -*-
from pathlib import Path
from datetime import date

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent
SAMPLES = ROOT / "samples"
SAMPLES.mkdir(parents=True, exist_ok=True)

# ---------------- Excel sample ----------------
wb = Workbook()
ws = wb.active
ws.title = "인증서데이터"
rows = [
    ["이름", "과정명", "소속", "수료일", "증서번호"],
    ["홍길동", "생성형 AI 실무과정", "플랫폼개발팀", date(2026, 8, 28), "ATA-2026-001"],
    ["김민지", "AI Agent 개발 실무", "기술교육팀", date(2026, 8, 28), "ATA-2026-002"],
    ["박서준", "클라우드 네이티브 실습", "SW개발팀", date(2026, 8, 29), "ATA-2026-003"],
]
for r in rows:
    ws.append(r)
header_fill = PatternFill("solid", fgColor="25364A")
for cell in ws[1]:
    cell.fill = header_fill
    cell.font = Font(bold=True, color="FFFFFF")
    cell.alignment = Alignment(horizontal="center", vertical="center")
for cell in ws["D"][1:]:
    cell.number_format = "yyyy-mm-dd"
for col, width in {"A": 14, "B": 28, "C": 20, "D": 14, "E": 18}.items():
    ws.column_dimensions[col].width = width
ws.freeze_panes = "A2"
wb.save(SAMPLES / "data.xlsx")

# ---------------- PowerPoint sample ----------------
prs = Presentation()
prs.slide_width = Inches(8.2677)
prs.slide_height = Inches(11.6929)
slide = prs.slides.add_slide(prs.slide_layouts[6])

NAVY = RGBColor(0x25, 0x36, 0x4A)
DARK = RGBColor(0x25, 0x28, 0x2C)
GRAY = RGBColor(0x73, 0x7A, 0x82)
LIGHT = RGBColor(0xD8, 0xDD, 0xE2)
PALE = RGBColor(0xF4, 0xF5, 0xF5)
WHITE = RGBColor(0xFC, 0xFC, 0xFB)

bg = slide.background.fill
bg.solid(); bg.fore_color.rgb = WHITE


def add_box(x, y, w, h, line_color, line_width=1.0, fill_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.color.rgb = line_color
    sh.line.width = Pt(line_width)
    if fill_color is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill_color
    return sh


def add_text(text, x, y, w, h, size, color=DARK, bold=False, align=PP_ALIGN.CENTER, font="맑은 고딕"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box

# borders
add_box(0.34, 0.34, 7.5877, 11.0129, NAVY, 1.4)
add_box(0.48, 0.48, 7.3077, 10.7329, LIGHT, 0.6)
add_text("증서번호  {증서번호}", 0.78, 0.72, 2.8, 0.3, 9, GRAY, False, PP_ALIGN.LEFT)
add_text("CERTIFICATE OF COMPLETION", 0.9, 1.44, 6.4677, 0.36, 10, GRAY, True, PP_ALIGN.CENTER, "Aptos")
add_text("수 료 증", 0.9, 1.86, 6.4677, 0.75, 30, NAVY, True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.72), Inches(2.74), Inches(2.82), Pt(1.5))
line.fill.solid(); line.fill.fore_color.rgb = NAVY; line.line.fill.background()

add_text("{이름}", 1.05, 3.36, 6.1677, 0.7, 24, DARK, True)
add_text("{소속}", 1.2, 4.06, 5.8677, 0.38, 11, GRAY)
add_text("위 사람은 첨단기술아카데미에서 실시한 아래 과정을\n성실히 이수하였으므로 이 증서를 수여합니다.", 1.15, 5.00, 5.9677, 0.9, 13.5, DARK)
add_box(1.37, 6.08, 5.5277, 1.02, LIGHT, 0.7, PALE)
add_text("{과정명}", 1.62, 6.31, 5.0277, 0.52, 16, NAVY, True)
add_text("{수료일}", 1.2, 8.01, 5.8677, 0.4, 12, DARK)
add_text("첨단기술아카데미", 1.2, 9.12, 5.8677, 0.5, 17, NAVY, True)
add_text("ADVANCED TECHNOLOGY ACADEMY", 1.2, 9.62, 5.8677, 0.28, 8.5, GRAY, False, PP_ALIGN.CENTER, "Aptos")

prs.save(SAMPLES / "인증서_{이름}_{과정명}.pptx")
print("samples created:", SAMPLES)
