from __future__ import annotations

import sys
from pathlib import Path

# `python scripts/...` 실행 시, 설치 없이 로컬 패키지를 import 가능하게 한다.
PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT))

import argparse

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from officekit.io_excel import generate_sample_excel
from officekit.utils import ensure_dir


def _create_word_template(path: Path) -> None:
    doc = Document()
    doc.add_heading("주문 요약서", level=1)
    doc.add_paragraph("아래 항목은 Excel 컬럼 값으로 자동 치환됩니다.")

    fields = [
        ("아이디", "{{ 아이디 }}"),
        ("이름", "{{ 이름 }}"),
        ("부서", "{{ 부서 }}"),
        ("분류", "{{ 분류 }}"),
        ("주문일자", "{{ 주문일자 }}"),
        ("이메일", "{{ 이메일 }}"),
        ("수량", "{{ 수량 }}"),
        ("단가", "{{ 단가 }}"),
        ("합계", "{{ 합계 }}"),
    ]

    table = doc.add_table(rows=len(fields), cols=2)
    table.style = "Table Grid"
    for idx, (label, value) in enumerate(fields):
        table.cell(idx, 0).text = label
        table.cell(idx, 1).text = value

    doc.save(str(path))


def _create_ppt_template(path: Path) -> None:
    prs = Presentation()

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    if cover.shapes.title is not None:
        cover.shapes.title.text = "엑셀 보고서"
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = "부제목"

    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    if table_slide.shapes.title is not None:
        table_slide.shapes.title.text = "데이터 표"
    table_area = table_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.2))
    table_area.name = "표_영역"
    table_area.text_frame.text = "표_영역"
    for p in table_area.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(20)

    chart_slide = prs.slides.add_slide(prs.slide_layouts[5])
    if chart_slide.shapes.title is not None:
        chart_slide.shapes.title.text = "분류별 합계"
    chart_area = chart_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.4))
    chart_area.name = "차트_영역"
    chart_area.text_frame.text = "차트_영역"
    for p in chart_area.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(20)

    prs.save(str(path))


def main() -> int:
    parser = argparse.ArgumentParser(description="한글 컬럼/식별자 기반 템플릿 세트를 생성합니다.")
    parser.add_argument("--templates-dir", default=str(PROJECT_ROOT / "templates"), help="템플릿 출력 디렉터리")
    parser.add_argument("--rows", type=int, default=30, help="한글 Excel 템플릿 샘플 행 수")
    args = parser.parse_args()

    out_dir = ensure_dir(args.templates_dir)

    excel_path = out_dir / "excel_template_ko.xlsx"
    word_path = out_dir / "word_template_ko.docx"
    ppt_path = out_dir / "ppt_template_ko.pptx"

    generate_sample_excel(excel_path, rows=args.rows, korean_mode=True)
    _create_word_template(word_path)
    _create_ppt_template(ppt_path)

    print("한글 템플릿 생성 완료:")
    print(f" - {excel_path}")
    print(f" - {word_path}")
    print(f" - {ppt_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
