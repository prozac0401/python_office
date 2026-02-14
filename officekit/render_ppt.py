from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

import matplotlib.pyplot as plt
from matplotlib import font_manager as fm

from .io_excel import read_excel_to_df
from .transform import add_total, summarize_by_group
from .utils import ensure_dir


def _find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if getattr(shape, "name", "") == name:
            return shape
    return None


def _find_shape_by_names(slide, names: tuple[str, ...]):
    for name in names:
        found = _find_shape_by_name(slide, name)
        if found is not None:
            return found
    return None


def _remove_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def _set_title(slide, title: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = title
        return
    # 대체: 첫 번째 텍스트 상자
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.text = title
            return


def _set_subtitle(slide, subtitle: str) -> None:
    # 일반적으로 placeholder[1]이 표지 슬라이드의 부제목이다.
    try:
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
            return
    except Exception:
        pass
    # 대체: 이름 기반 도형 탐색
    shp = _find_shape_by_names(slide, ("SUBTITLE", "부제목"))
    if shp and shp.has_text_frame:
        shp.text_frame.text = subtitle


def _get_or_create_slide(prs: Presentation, index: int, layout_index: int):
    if len(prs.slides) > index:
        return prs.slides[index]
    return prs.slides.add_slide(prs.slide_layouts[layout_index])


def _get_region_from_named_shape(slide, shape_names: tuple[str, ...], default_region):
    shp = _find_shape_by_names(slide, shape_names)
    if shp is None:
        return default_region, None
    region = (shp.left, shp.top, shp.width, shp.height)
    return region, shp


def _insert_table(slide, df: pd.DataFrame, region, max_rows: int = 12) -> None:
    left, top, width, height = region
    view = df.head(max_rows)

    rows = len(view) + 1
    cols = len(view.columns)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 헤더
    for j, col in enumerate(view.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

    # 본문
    for i in range(len(view)):
        for j, col in enumerate(view.columns):
            val = view.iloc[i, j]
            table.cell(i + 1, j).text = "" if pd.isna(val) else str(val)

    # 폰트 크기/정렬
    for r in range(rows):
        for c in range(cols):
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                if p.font.size is None:
                    p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER


def _make_bar_chart_png(summary_df: pd.DataFrame, category_col: str, value_col: str, out_path: Path, title: str) -> Path:
    available_fonts = {f.name for f in fm.fontManager.ttflist}
    for name in ("Malgun Gothic", "맑은 고딕", "NanumGothic", "Noto Sans CJK KR", "AppleGothic"):
        if name in available_fonts:
            plt.rcParams["font.family"] = name
            break
    plt.rcParams["axes.unicode_minus"] = False

    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig = plt.figure(figsize=(10, 5.5))
    ax = fig.add_subplot(111)

    x = summary_df[category_col].astype(str).tolist()
    y = summary_df[value_col].tolist()
    ax.bar(x, y)
    ax.set_title(title)
    ax.set_xlabel(category_col)
    ax.set_ylabel(value_col)
    ax.tick_params(axis="x", rotation=30)

    fig.tight_layout()
    fig.savefig(out_path, dpi=200)
    plt.close(fig)
    return out_path


def _insert_picture(slide, img_path: Path, region) -> None:
    left, top, width, height = region
    slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)


def create_ppt_from_excel(
    input_xlsx: str | Path,
    output_pptx: str | Path,
    sheet_name: int | str = 0,
    title: str = "엑셀 보고서",
    group_col: str = "Category",
    qty_col: str = "Qty",
    unit_price_col: str = "UnitPrice",
    template_pptx: str | Path | None = None,
) -> Path:
    """Excel 파일에서 데이터 기반 PPT 보고서를 생성한다(COM 미사용).

    `template_pptx`를 지정하면 해당 파일을 기반으로 생성한다.
    프로젝트 기본 템플릿은 다음 구조를 가정한다.
    - 1번 슬라이드: 표지
    - 2번 슬라이드: 표 슬라이드, 이름이 "TABLE_AREA" 또는 "표_영역"인 사각형 도형 포함
    - 3번 슬라이드: 차트 슬라이드, 이름이 "CHART_AREA" 또는 "차트_영역"인 사각형 도형 포함

    위 슬라이드/도형이 없으면 단순 레이아웃으로 대체한다.
    """
    df = read_excel_to_df(input_xlsx, sheet_name=sheet_name)
    df2 = add_total(df, qty_col=qty_col, unit_price_col=unit_price_col, out_col="Total")
    summary = summarize_by_group(df2, group_col=group_col, value_col="Total", agg="sum")
    summary_group_col = str(summary.columns[0]) if len(summary.columns) > 0 else group_col
    summary_value_col = str(summary.columns[1]) if len(summary.columns) > 1 else "Total"

    # 템플릿이 지정되면 템플릿을 로드한다.
    if template_pptx:
        tpl_path = Path(template_pptx)
        if not tpl_path.exists():
            raise FileNotFoundError(f"템플릿 파일을 찾을 수 없습니다: {tpl_path}")
        prs = Presentation(str(tpl_path))
    else:
        prs = Presentation()

    # ---- 1번 슬라이드(표지)
    cover = _get_or_create_slide(prs, 0, 0)
    _set_title(cover, title)
    _set_subtitle(cover, f"원본: {Path(input_xlsx).name}")

    # ---- 2번 슬라이드(표)
    table_slide = _get_or_create_slide(prs, 1, 5)
    _set_title(table_slide, "데이터 미리보기(상위 행)")
    default_table_region = (Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.2))
    region, placeholder = _get_region_from_named_shape(table_slide, ("TABLE_AREA", "표_영역"), default_table_region)
    if placeholder is not None:
        _remove_shape(placeholder)
    _insert_table(table_slide, df2, region, max_rows=12)

    # ---- 3번 슬라이드(차트)
    chart_slide = _get_or_create_slide(prs, 2, 5)
    _set_title(chart_slide, f"{summary_group_col}별 {summary_value_col}")
    tmp_dir = ensure_dir(Path(output_pptx).parent / "_tmp")
    chart_png = _make_bar_chart_png(
        summary,
        category_col=summary_group_col,
        value_col=summary_value_col,
        out_path=tmp_dir / "total_by_category.png",
        title=f"{summary_group_col}별 {summary_value_col}",
    )
    default_chart_region = (Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.4))
    region2, placeholder2 = _get_region_from_named_shape(chart_slide, ("CHART_AREA", "차트_영역"), default_chart_region)
    if placeholder2 is not None:
        _remove_shape(placeholder2)
    _insert_picture(chart_slide, chart_png, region2)

    out = Path(output_pptx)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
